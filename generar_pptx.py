from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def add_image_placeholder(slide, left, top, width, height, text):
    """Añade un recuadro gris simulando el hueco para la imagen"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = text
    return shape

def crear_presentacion():
    prs = Presentation()

    # Slide 1: Objective
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Layout de Título y Contenido
    slide.shapes.title.text = "Objective: Timing Side-Channel in RAG Systems"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Main Goal: To demonstrate that an external attacker can infer the presence, frequency, and semantic distribution of private terms within a RAG system."
    body.add_paragraph().text = "Attack Vector: Exploiting query latency variances in the underlying Vector Database (ChromaDB) based on structural graph traversals."
    body.add_paragraph().text = "Threat Model: Zero-knowledge attacker. No direct database access, no metadata extraction, relying purely on measuring response times."

    # Slide 2: Experimental Setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Experimental Setup: Models & Dataset"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Corpus: emozilla/pg19 (Project Gutenberg Books), fixed 1000-character chunks."
    body.add_paragraph().text = "Vector Database: ChromaDB utilizing default HNSW indexing."
    body.add_paragraph().text = "Embeddings: all-MiniLM-L6-v2 (Dense vector mapping)."
    body.add_paragraph().text = "LLMs Tested (Generation Phase):"
    p = body.add_paragraph()
    p.text = "- Large Model: Qwen-2.5-7B-Instruct (High parameter count, realistic generation)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "- Tiny Model: SmolLM2-135M-Instruct (Minimal parameter count, isolated inference)."
    p.level = 1

    # Slide 3: Direct RAG Results
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Layout de solo título
    slide.shapes.title.text = "Baseline Results: Direct Database Queries"
    add_image_placeholder(slide, Inches(0.5), Inches(2), Inches(4), Inches(3.5), "[INSERT IMAGE: Whisker Plot k=5 (Flat/Uniform)]")
    add_image_placeholder(slide, Inches(5.5), Inches(2), Inches(4), Inches(3.5), "[INSERT IMAGE: Whisker Plot k=50 (Step-pattern)]")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    txBox.text_frame.text = "Conclusion: With big K, the search algorithm traverses deeper into the HNSW graph, exposing the semantic density and distribution of the private data."

    # Slide 4: Results with Large Model
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "The Noise Problem: RAG Pipeline with Qwen-2.5-7B"
    add_image_placeholder(slide, Inches(0.5), Inches(2), Inches(4), Inches(3.5), "[INSERT IMAGE: Whisker Plot Qwen k=5]")
    add_image_placeholder(slide, Inches(5.5), Inches(2), Inches(4), Inches(3.5), "[INSERT IMAGE: Whisker Plot Qwen k=50]")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    txBox.text_frame.text = "Conclusion: The time taken by Qwen makes the database latencies completely indistinguishable."

    # Slide 5: Results with Tiny Model
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Isolating the Leak: RAG Pipeline with SmolLM2-135M"
    add_image_placeholder(slide, Inches(0.5), Inches(2), Inches(4), Inches(3.5), "[INSERT IMAGE: Whisker Plot SmolLM2 k=5]")
    add_image_placeholder(slide, Inches(5.5), Inches(2), Inches(4), Inches(3.5), "[INSERT IMAGE: Whisker Plot SmolLM2 k=50]")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    txBox.text_frame.text = "Conclusion: By minimizing generation overhead, we see exactly the same structural latency leaks as directly with the database."

    # Slide 6: The Turn Inference Tool
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Next Steps: The Turn Inference Tool"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "The Ultimate Goal: Moving from empirical observation to a fully functional statistical attack."
    body.add_paragraph().text = "The Inference Tool: Building a classifier that ingests timing data and outputs confidence scores regarding the internal database contents."
    body.add_paragraph().text = "Objective: To precisely calculate how many query rounds (turns) an attacker needs to execute to statistically guarantee whether a specific hidden term exists or not in the RAG system."

    # Guardar
    prs.save('TFG_RAG_Presentation.pptx')
    print("[+] Presentación generada con éxito: TFG_RAG_Presentation.pptx")

if __name__ == "__main__":
    crear_presentacion()